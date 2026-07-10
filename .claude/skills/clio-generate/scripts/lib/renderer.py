from __future__ import annotations

import base64
import copy
import os
import tempfile
from pathlib import Path
from typing import Dict
from typing import List
from typing import Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

from lib.schemas import ContentItem
from lib.schemas import ContentType
from lib.schemas import SlideContent
from lib.templates import SVNProposalTemplate
from lib.templates.base import ContentKind, ShapeRole, SlideRoleConfig
from lib.profile_schema import ProjectProfile
from lib.shape_discovery import ShapeDiscovery



class PPTXRenderer:
    """Renderer for filling content into existing template slides"""

    # OOXML unit: 914400 EMU = 1 inch.  All other EMU constants derive from this.
    _EMU_PER_INCH: int = 914400

    # Rows are compressed to fit the safe area; this floor keeps text legible.
    _MIN_ROW_HEIGHT_EMU: int = int(0.25 * _EMU_PER_INCH)

    # Minimum height for data rows to ensure content is readable.
    # At 8pt Noto Sans JP, one line ≈ 0.15"; two lines + padding ≈ 0.35".
    _MIN_DATA_ROW_HEIGHT_EMU: int = int(0.35 * _EMU_PER_INCH)

    def __init__(self):
        self.template_config = None
        self._output_dir: Optional[Path] = None
        # Safe content area bottom — computed dynamically from slide height in render_from_profile.
        # Fallback: 5" (suits a 5.625" widescreen slide with footer at ~5.25").
        self._SVN_SAFE_BOTTOM_EMU: int = int(5.0 * self._EMU_PER_INCH)
        # Actual slide width — set from prs.slide_width in render_from_profile.
        self._slide_width_emu: int = 0

    def render(
        self,
        slides: List[SlideContent],
        template: str,
        output_name: str,
        output_dir: Optional[str] = None,
    ) -> Path:
        """
        Fill content into existing template slides

        Args:
            slides: List of slide content to fill
            template: Template PPTX file path (e.g., 'SVN Proposal Menu.pptx')
            output_name: Output file name (without .pptx extension)
            output_dir: Deprecated/ignored. Output directory is resolved from env `SLIDE_GENERATOR__OUTPUTS_PATH`.

        Returns:
            Path to generated PPTX file
        """
        print(f'Filling content into template slides: {output_name}')

        # Resolve output_dir once so image path resolution can use it
        self._output_dir = self._resolve_output_dir(output_dir)

        # Auto-detect template config from template name
        self._detect_template_config(template)

        # Load template presentation
        prs = self._load_template(template)

        # Fill content into existing slides
        self._fill_existing_slides(prs, slides)

        # Save presentation
        output_file = self._save_presentation(prs, output_name, output_dir)

        return output_file

    @staticmethod
    def _resolve_output_dir(output_dir: Optional[str]) -> Path:
        """Resolve output directory path.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        if output_dir:
            return Path(output_dir)
        env_outputs = os.getenv('SLIDE_GENERATOR__OUTPUTS_PATH')
        if env_outputs:
            return Path(env_outputs).expanduser()
        return Path.cwd() / 'outputs'

    def _detect_template_config(self, template: str):
        """Auto-detect and load template config from template name"""
        if 'svn' in template.lower():
            self.template_config = SVNProposalTemplate()
            print(f'Auto-detected SVN template config: {len(self.template_config.slide_configs)} slides configured')
        else:
            self.template_config = None
            print('No specific template config detected, using default rendering')

    def _load_template(self, template: str) -> Presentation:
        """Load template presentation — supports absolute paths and relative names."""
        template_path = Path(template)
        if template_path.is_absolute() and template_path.exists():
            prs = Presentation(str(template_path))
            print(f'Loaded template: {template_path.name} ({len(prs.slides)} slides)')
            return prs
        rel_path = Path(__file__).parent / 'templates' / template
        if rel_path.exists():
            prs = Presentation(str(rel_path))
            print(f'Loaded template: {template} from {rel_path} ({len(prs.slides)} slides)')
            return prs
        print(f'Template not found: {template} at {template_path} or {rel_path}')
        return Presentation()

    def _save_presentation(self, prs: Presentation, output_name: str, output_dir: Optional[str] = None) -> Path:
        """Save presentation to output directory.

        Priority: explicit output_dir > env SLIDE_GENERATOR__OUTPUTS_PATH > CWD/outputs
        """
        output_path = self._output_dir or self._resolve_output_dir(output_dir)

        try:
            output_path.mkdir(parents=True, exist_ok=True)
        except PermissionError as e:
            # Make error clearer for API layer
            raise PermissionError(f'Permission denied creating outputs directory: {output_path}') from e
        output_file = output_path / f'{output_name}.pptx'

        prs.save(str(output_file))
        print(f'Saved PPTX to: {output_file}')

        return output_file

    def _fill_existing_slides(self, prs: Presentation, slides_content: List[SlideContent]):
        """Fill content into existing slides"""
        print(f'Filling content into {len(slides_content)} existing slides')

        # Sort by slide number so that the offset tracking works correctly when
        # overflow slides are inserted in the middle of the presentation.
        sorted_contents = sorted(slides_content, key=lambda s: s.slide_number)

        # Tracks how many extra slides have been inserted before the current one.
        slide_offset = 0

        for slide_content in sorted_contents:
            slide_num = slide_content.slide_number

            # Check if slide is protected
            if self.template_config and not self.template_config.should_fill_slide(slide_num):
                print(f'Slide {slide_num} is protected - skipping')
                continue

            # Adjusted index accounts for any extra slides inserted before this slide.
            adjusted_index = slide_num - 1 + slide_offset

            # Check if slide exists
            if adjusted_index < len(prs.slides):
                slide = prs.slides[adjusted_index]
                print(f'Filling content into slide {slide_num} (template index: {adjusted_index})')

                # Fill content; if a table overflows it will insert extra slides and
                # return the count so we can update the offset.
                extra = self._fill_with_overflow_handling(prs, slide, slide_content, adjusted_index)
                slide_offset += extra
            else:
                print(
                    f'Slide {slide_num} not found in template '
                    f'(total slides: {len(prs.slides)}, adjusted index: {adjusted_index})',
                )

    # ------------------------------------------------------------------
    # Table-overflow helpers
    # ------------------------------------------------------------------

    def _fill_with_overflow_handling(
        self,
        prs: Presentation,
        slide,
        slide_content: SlideContent,
        slide_index: int,
    ) -> int:
        """Fill a slide, splitting into multiple slides when a table is too tall
        to fit within the safe content area above the footer band.

        Returns the number of *extra* slides inserted (0 when no overflow).
        """
        shape_contents = getattr(slide_content, 'shape_contents', {}) or {}

        # Find the first table-type shape content for this slide
        table_key, table_data = self._find_table_content(slide_content, shape_contents)

        if not table_key:
            # No table present – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        # Legacy path: shape position is unknown here, so estimate from safe area.
        # Assume content starts ~1.5" below the top (title band); minimum row
        # height gives the theoretical maximum rows that could ever fit.
        data_area_est = self._SVN_SAFE_BOTTOM_EMU - int(1.5 * self._EMU_PER_INCH)
        legacy_max_rows = max(1, data_area_est // self._MIN_ROW_HEIGHT_EMU)
        chunks = self._split_table_data(table_data, legacy_max_rows)

        if len(chunks) <= 1:
            # Fits on one slide – fill normally
            self._fill_slide_content(slide, slide_content)
            return 0

        print(
            f'  - Slide {slide_content.slide_number}: table split '
            f'({len(chunks)} page(s), up to {legacy_max_rows} rows/page)',
        )

        # Fill the original slide with the first chunk
        first_content = self._build_chunked_slide_content(
            slide_content, table_key, chunks[0], is_first=True,
        )
        self._fill_slide_content(slide, first_content)

        # Create one continuation slide per remaining chunk
        extra_slides = 0
        current_index = slide_index
        for i, chunk in enumerate(chunks[1:], 1):
            new_slide = self._duplicate_slide(prs, current_index)
            extra_slides += 1
            current_index += 1

            cont_content = self._build_chunked_slide_content(
                slide_content, table_key, chunk, is_first=False,
            )
            self._fill_slide_content(new_slide, cont_content)
            print(
                f'  - Added continuation slide {i}/{len(chunks) - 1} '
                f'at template position {current_index + 1}',
            )

        return extra_slides

    def _find_table_content(
        self,
        slide_content: SlideContent,
        shape_contents: dict,
    ):
        """Return (content_key, markdown_table_string) for the first table-type
        shape found in this slide's config.  Returns (None, None) if none found.
        """
        config = (
            self.template_config.get_slide_config(slide_content.slide_number)
            if self.template_config
            else None
        )
        if not config or not config.shape_targets:
            return None, None

        for target in config.shape_targets:
            content = shape_contents.get(target.content_key, '')
            if content and self._is_table_data(content):
                return target.content_key, content

        return None, None

    def _split_table_data(self, table_data: str, max_rows: int) -> List[str]:
        """Split a markdown table string into chunks of at most *max_rows* data rows.

        Each chunk keeps the original header row(s) and the separator row so that
        each chunk is a valid, self-contained markdown table.

        Returns a list of markdown table strings (length 1 when no split needed).
        """
        lines = [line for line in table_data.strip().split('\n') if line.strip()]
        if not lines:
            return [table_data]

        # Separate header / separator lines from data lines
        header_lines: List[str] = []
        data_lines: List[str] = []
        separator_found = False

        for line in lines:
            stripped = line.strip()
            is_separator = (
                not separator_found
                and '|' in stripped
                and all(c in '|-: \t' for c in stripped)
            )
            if is_separator:
                header_lines.append(line)
                separator_found = True
            elif not separator_found:
                header_lines.append(line)
            else:
                data_lines.append(line)

        if len(data_lines) <= max_rows:
            return [table_data]

        header_block = '\n'.join(header_lines)
        chunks: List[str] = []
        for i in range(0, len(data_lines), max_rows):
            row_block = '\n'.join(data_lines[i : i + max_rows])
            chunks.append(f'{header_block}\n{row_block}')

        return chunks

    def _build_chunked_slide_content(
        self,
        original: SlideContent,
        table_key: str,
        table_chunk: str,
        is_first: bool,
    ) -> SlideContent:
        """Return a new SlideContent with *table_key* replaced by *table_chunk*.

        For continuation slides (``is_first=False``) all non-table text shapes
        are cleared so description text is not repeated across slides.
        """
        new_shape_contents = dict(getattr(original, 'shape_contents', {}) or {})
        new_shape_contents[table_key] = table_chunk

        if not is_first:
            config = (
                self.template_config.get_slide_config(original.slide_number)
                if self.template_config
                else None
            )
            if config and config.shape_targets:
                for target in config.shape_targets:
                    key = target.content_key
                    if key != table_key and key in new_shape_contents:
                        if not self._is_table_data(new_shape_contents[key]):
                            new_shape_contents[key] = ''  # Clear description text

        return SlideContent(
            slide_number=original.slide_number,
            layout=original.layout,
            title=original.title,
            subtitle=original.subtitle,
            content=original.content,
            shape_contents=new_shape_contents,
        )

    def _duplicate_slide(self, prs: Presentation, slide_index: int):
        """Duplicate the slide at *slide_index* and insert the copy immediately
        after it in the presentation.

        Returns the new slide object (at position ``slide_index + 1``).
        """
        source_slide = prs.slides[slide_index]

        # 1. Add a new blank slide with the same layout – this is appended at the end.
        new_slide = prs.slides.add_slide(source_slide.slide_layout)

        # 2. Replace the new slide's shape tree with a deep copy of the source's.
        source_spTree = source_slide.shapes._spTree
        new_spTree = new_slide.shapes._spTree

        for child in list(new_spTree):
            new_spTree.remove(child)
        for child in source_spTree:
            new_spTree.append(copy.deepcopy(child))

        # 3. Move the new slide from the end to slide_index + 1 by reordering
        #    the <p:sldIdLst> in the presentation XML.
        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            presentation_elem = prs.part._element
            sldIdLst = presentation_elem.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                sld_id_elems = list(sldIdLst)
                if sld_id_elems:
                    new_sldId = sld_id_elems[-1]          # just appended → at end
                    source_sldId = sld_id_elems[slide_index]  # source position
                    sldIdLst.remove(new_sldId)
                    source_sldId.addnext(new_sldId)
                    print(
                        f'  - Duplicated slide at index {slide_index}, '
                        f'inserted at index {slide_index + 1}',
                    )
        except Exception as e:
            print(f'  - Error repositioning duplicated slide: {e}')

        return prs.slides[slide_index + 1]

    # ------------------------------------------------------------------

    def _fill_slide_content(self, slide, slide_content: SlideContent):
        """Fill content into an existing slide"""
        slide_num = slide_content.slide_number
        config = self.template_config.get_slide_config(slide_num) if self.template_config else None

        # Strategy 1: Shape targeting (if shape_contents exists)
        if self._has_shape_contents(slide_content) and config and config.shape_targets:
            print(f'  - Using shape targeting for slide {slide_num}')
            self._fill_with_shape_targets(slide, slide_content, config)
            return

        # Strategy 2: Default filling
        self._fill_slide_default(slide, slide_content)

    def _has_shape_contents(self, slide_content: SlideContent) -> bool:
        """Check if slide_content has shape_contents attribute"""
        return bool(
            hasattr(slide_content, '__dict__') and
            'shape_contents' in slide_content.__dict__ and
            slide_content.shape_contents,
        )

    def _fill_with_shape_targets(self, slide, slide_content, config):
        """Fill content using shape targets from config"""
        shape_contents = getattr(slide_content, 'shape_contents', {})

        if not shape_contents:
            print('  - No shape_contents found in slide_content')
            return

        print(f'  - Found {len(shape_contents)} shape contents to fill')

        # Update title if exists
        self._update_slide_title(slide, slide_content.title)

        # Fill each shape target
        for target in config.shape_targets:
            self._fill_single_shape_target(slide, target, shape_contents, config)

    def _update_slide_title(self, slide, title: Optional[str]):
        """Update slide title if exists"""
        if slide.shapes.title and title:
            slide.shapes.title.text = title
            print(f'  - Updated title: {title}')

    def _fill_single_shape_target(self, slide, target, shape_contents: Dict[str, str], config=None):
        """Fill a single shape target with content"""
        content_key = target.content_key

        # Check if content exists for this key
        if content_key not in shape_contents:
            print(f"  - Content key '{content_key}' not found in markdown")
            return

        content = shape_contents[content_key]

        # Find shape by index or name
        shape = self._find_shape(slide, target)

        if not shape:
            print(f'  - Shape not found for target: {content_key}')
            return

        # Detect content type and fill accordingly
        if self._is_image_path(content):
            # Content is an image path
            self._replace_shape_with_image(slide, shape, content, target)
        elif target.fill_cols or self._is_table_data(content):
            # fill_cols always implies a table target (even single-row);
            # otherwise fall back to markdown table detection.
            self._fill_table(shape, content, target)
        elif shape.has_text_frame:
            # Content is text
            shape.text_frame.text = content
            self._apply_text_font(shape.text_frame)
            
            # Apply 10pt line spacing for slide 4
            if config and config.slide_number == 4:
                self._apply_line_spacing(shape.text_frame, Pt(10))
                print(f"  - Applied 10pt line spacing to slide 4 shape '{content_key}'")
            
            print(f"  Filled '{content_key}' into {shape.name} ({len(content)} chars)")
        else:
            print(f'  - Shape {shape.name} has no text frame and content is not an image or table')

    def _is_image_path(self, content: str) -> bool:
        """Check if content is an image file path or base64 data URI"""
        content = content.strip()
        # Check base64 data URI: data:image/...;base64,...
        if content.startswith('data:image/'):
            return True
        # Check if it's a file path with image extension
        image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg']
        return any(content.lower().endswith(ext) for ext in image_extensions)

    def _decode_base64_image(self, data_uri: str) -> Optional[str]:
        """Decode a base64 data URI to a temporary file and return the file path.

        Supports format: data:image/<ext>;base64,<data>

        Returns:
            Absolute path to the temporary file, or None on error.
        """
        try:
            # Parse header: data:image/png;base64,<data>
            header, encoded = data_uri.split(',', 1)
            # Extract mime type, e.g. image/png
            mime_part = header.split(':')[1].split(';')[0]  # e.g. "image/png"
            ext = '.' + mime_part.split('/')[1]  # e.g. ".png"
            # Handle special cases
            if ext == '.jpeg':
                ext = '.jpg'
            elif ext == '.svg+xml':
                ext = '.svg'

            image_bytes = base64.b64decode(encoded)
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp.write(image_bytes)
            tmp.flush()
            tmp.close()
            print(f'  - Decoded base64 image to temp file: {tmp.name} ({len(image_bytes)} bytes)')
            return tmp.name
        except Exception as e:
            print(f'  - Failed to decode base64 image: {e}')
            return None

    def _replace_shape_with_image(self, slide, old_shape, image_path: str, target):
        """Replace a shape with an image (supports file path or base64 data URI)"""
        tmp_file: Optional[str] = None
        try:
            # Get shape position and size
            left = old_shape.left
            top = old_shape.top
            width = old_shape.width
            height = old_shape.height

            image_path = image_path.strip()

            # Handle base64 data URI
            if image_path.startswith('data:image/'):
                tmp_file = self._decode_base64_image(image_path)
                if not tmp_file:
                    print(f"  - Cannot decode base64 image for shape '{target.content_key}'")
                    return
                image_path = tmp_file
            else:
                # Resolve file path
                if not Path(image_path).is_absolute():
                    rel_path = Path(image_path)
                    if rel_path.exists():
                        image_path = str(rel_path)
                    elif self._output_dir:
                        resolved = self._output_dir / image_path
                        if resolved.exists():
                            image_path = str(resolved)
                    else:
                        outputs_path = Path('outputs') / image_path
                        if outputs_path.exists():
                            image_path = str(outputs_path)

                # Check if image exists
                if not Path(image_path).exists():
                    print(f'  - Image not found: {image_path}')
                    return

            # Add new image at same position (overlays the placeholder shape)
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    left, top,
                    width=width,
                    height=height,
                )
                print(f"  Replaced shape '{target.content_key}' with image: {Path(image_path).name}")
            except Exception as e:
                print(f'  - Error adding image: {e}')

        except Exception as e:
            print(f'  - Error replacing shape with image: {e}')
        finally:
            # Clean up temp file if created from base64
            if tmp_file and Path(tmp_file).exists():
                try:
                    Path(tmp_file).unlink()
                    print(f'  - Cleaned up temp file: {tmp_file}')
                except Exception:
                    pass

    def _is_table_data(self, content: str) -> bool:
        """Check if content is table data (markdown table format).

        Accepts two formats:
        1. Standard markdown table with header + separator:
               | col1 | col2 |
               |------|------|
               | val1 | val2 |
        2. No-header row-only table (used for fill_cols targeting):
               | val1 |
               | val2 |
               | val3 |
        """
        content = content.strip()
        lines = [line.strip() for line in content.split('\n') if line.strip()]

        if len(lines) < 2:
            return False

        # All lines must contain at least one pipe
        if not all('|' in line for line in lines):
            return False

        # Format 1: standard table with separator line
        for line in lines[1:3]:
            if all(c in '|-: \t' for c in line):
                return True

        # Format 2: no-separator — every line starts and ends with | (row-only table)
        if all(line.startswith('|') and line.endswith('|') for line in lines):
            return True

        return False

    def _fill_table(self, shape, table_data: str, target):
        """Fill an existing table shape with data from markdown table format"""
        try:
            # Check if shape has table
            if not shape.has_table:
                print(f"  - Shape '{target.content_key}' is not a table")
                return

            table = shape.table

            # Parse markdown table
            lines = [line.strip() for line in table_data.strip().split('\n') if line.strip()]

            # Remove separator line (e.g., |---|---|)
            table_lines = []
            for line in lines:
                # Skip separator lines
                if all(c in '|-: \t' for c in line):
                    continue
                table_lines.append(line)

            if not table_lines:
                print('  - No valid table data found')
                return

            # Parse cells from each line
            rows_data = []
            for line in table_lines:
                # Split by | and clean up
                cells = [cell.strip() for cell in line.split('|')]
                # Remove empty first/last cells (from leading/trailing |)
                cells = [cell for cell in cells if cell]
                if cells:
                    rows_data.append(cells)

            if not rows_data:
                print('  - No valid rows found in table data')
                return

            # Get current table dimensions
            current_rows = len(table.rows)
            current_cols = len(table.columns)
            needed_rows = len(rows_data)
            needed_cols = len(rows_data[0]) if rows_data else 0

            # Adjust table size if needed
            if needed_rows > current_rows:
                # Add rows if we need more
                rows_to_add = needed_rows - current_rows
                print(f'  - Adding {rows_to_add} rows to table (current: {current_rows}, needed: {needed_rows})')
                self._add_table_rows(table, rows_to_add)
                # After adding rows, update current_rows to reflect the new count.
                # Note: python-pptx's table.rows collection is live, so this should
                # reflect the new size immediately, but we recalculate for safety.
                current_rows = len(table.rows)
            elif needed_rows < current_rows:
                # Fewer rows needed - remove the extra rows from the XML so they
                # don't appear as empty rows in the rendered slide.
                extra_rows = current_rows - needed_rows
                print(
                    '  - Table has %s rows but only %s rows of data '
                    '(removing %s unused rows)',
                    current_rows,
                    needed_rows,
                    extra_rows,
                )
                self._remove_table_rows(table, needed_rows)
                current_rows = len(table.rows)

            if needed_cols > current_cols:
                print(f'  - Data has {needed_cols} columns but table only has {current_cols} columns')
                needed_cols = current_cols  # Limit to available columns

            # Fill table cells
            filled_count = 0
            for row_idx, row_data in enumerate(rows_data):
                # After adding rows via XML, we can fill up to needed_rows
                # Use try-except to handle any edge cases where table structure might not match
                try:
                    # If fill_cols is specified, map row_data positionally onto those columns;
                    # columns not listed are left untouched (preserving template content).
                    if target.fill_cols:
                        col_pairs = list(zip(target.fill_cols, row_data))
                    else:
                        col_pairs = list(enumerate(row_data))

                    for col_idx, cell_value in col_pairs:
                        if col_idx >= current_cols:
                            break  # Skip extra columns in data

                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_value
                        self._apply_table_cell_font(cell)
                        filled_count += 1
                except (IndexError, AttributeError) as e:
                    print(f'  - Cannot fill row {row_idx}: {e}')
                    break  # Stop filling if we hit an error

            print(f"  Filled '{target.content_key}' table: {len(rows_data)} rows x {needed_cols} cols ({filled_count} cells)")
            self._normalize_table_col_widths(table)
            self._apply_content_row_heights(table)
            self._fit_table_to_bounds(shape)

        except Exception as e:
            print(f'  - Error filling table: {e}')

    def _add_table_rows(self, table, num_rows: int):
        """Add rows to an existing table by manipulating XML

        python-pptx uses lxml internally, so we can manipulate the XML directly
        to add new rows by cloning the structure of the last row.
        """
        try:
            from lxml import etree

            # Get the table XML element (this is an lxml element)
            tbl = table._tbl

            # Get the number of existing rows
            current_row_count = len(table.rows)

            # Get the last row to use as template for new rows
            if current_row_count == 0:
                print('  - Cannot add rows: table has no existing rows to use as template')
                return

            # Use positive index instead of negative index (table.rows doesn't support negative indexing)
            last_row = table.rows[current_row_count - 1]
            last_row_xml = last_row._tr

            # Define namespace for XPath queries
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # Add new rows by cloning the last row
            for i in range(num_rows):
                # Create a deep copy of the last row by serializing and parsing
                new_row_xml = etree.fromstring(etree.tostring(last_row_xml))

                # Clear text content in all cells of the new row
                for tc in new_row_xml.xpath('.//a:tc', namespaces=ns):
                    txBody = tc.find('.//a:txBody', namespaces=ns)
                    if txBody is not None:
                        # Get all paragraphs
                        paragraphs = txBody.findall('.//a:p', namespaces=ns)
                        if paragraphs:
                            # Clear text in first paragraph
                            first_p = paragraphs[0]
                            for t_elem in first_p.xpath('.//a:t', namespaces=ns):
                                t_elem.text = ''
                            # Remove other paragraphs (keep structure but remove extra content)
                            for p in paragraphs[1:]:
                                txBody.remove(p)

                # Append the new row to the table
                tbl.append(new_row_xml)

            print(f'  - Successfully added {num_rows} rows to table')

        except ImportError:
            print('  - lxml not available. python-pptx requires lxml, please install it: pip install lxml')
        except Exception as e:
            print(f'  - Error adding rows to table: {e}')
            import traceback
            print(f'  - Traceback: {traceback.format_exc()}')

    def _remove_table_rows(self, table, keep_rows: int):
        """Remove all rows after index *keep_rows* from the table by manipulating XML.

        Args:
            table: python-pptx Table object
            keep_rows: Number of rows to keep (rows at index 0 … keep_rows-1 are retained).
        """
        try:
            from lxml import etree  # noqa: F401 (just to confirm availability)

            tbl = table._tbl
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            tr_elements = tbl.findall('a:tr', namespaces=ns)

            rows_to_remove = tr_elements[keep_rows:]
            for tr in rows_to_remove:
                tbl.remove(tr)

            print(
                f'  - Removed {len(rows_to_remove)} extra row(s); '
                f'table now has {keep_rows} row(s)',
            )
        except Exception as e:
            print(f'  - Error removing table rows: {e}')

    def _apply_content_row_heights(self, table) -> None:
        """Bump data-row heights so each row is tall enough for its actual text content.

        The template may have very compact rows (e.g. 0.262") designed for short
        placeholder text.  When real content is longer and word-wraps, some renderers
        (LibreOffice, Google Slides preview) auto-expand rows — causing the table to
        grow past the footer even though the python-pptx shape is technically in-bounds.

        This method estimates the required row height from text length + column width,
        then sets each data row to max(template_height, content_height, MIN_DATA_ROW).
        After this call, _fit_table_to_bounds sees the true height and can compress
        or trigger a split correctly.
        """
        try:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            emu = self._EMU_PER_INCH

            cols = list(table.columns)
            col_widths_inch = [max(c.width / emu, 0.5) for c in cols] if cols else []

            # Font metrics at 8pt Noto Sans JP
            font_pt = 8
            cjk_w_inch = font_pt / 72          # ~0.111" per CJK char (1 em)
            latin_w_inch = font_pt * 0.6 / 72   # ~0.067" per Latin char
            line_h_emu = int(font_pt * 1.6 / 72 * emu)   # ~0.178" line height with spacing
            padding_emu = int(0.07 * emu)        # top + bottom cell padding ~0.07"

            tbl = table._tbl
            tr_elements = tbl.findall('a:tr', namespaces=ns)
            changed = 0

            for idx, tr in enumerate(tr_elements):
                if idx == 0:
                    continue  # skip header row

                current_h = int(tr.get('h', '0'))
                min_h = max(current_h, self._MIN_DATA_ROW_HEIGHT_EMU)

                # Estimate max lines needed across all cells in this row
                max_lines = 1
                for j, tc in enumerate(tr.findall('a:tc', namespaces=ns)):
                    if j >= len(col_widths_inch):
                        break
                    # Extract cell text
                    texts = tc.findall('.//a:t', namespaces=ns)
                    text = ''.join(t.text or '' for t in texts).strip()
                    if not text:
                        continue
                    col_w = col_widths_inch[j]
                    # Estimate text width
                    text_width = sum(
                        cjk_w_inch if ('　' <= ch <= '鿿' or '＀' <= ch <= '￯')
                        else latin_w_inch
                        for ch in text
                    )
                    lines = max(1, int(text_width / col_w) + 1)
                    max_lines = max(max_lines, lines)

                content_h = int(max_lines * line_h_emu + padding_emu)
                new_h = max(min_h, content_h)

                if new_h != current_h:
                    tr.set('h', str(new_h))
                    changed += 1

            if changed:
                print(f'  - Content-adjusted {changed} row height(s) '
                      f'(min {self._MIN_DATA_ROW_HEIGHT_EMU/emu:.3f}")')
        except Exception as e:
            print(f'  - _apply_content_row_heights error: {e}')

    def _max_rows_for_shape(self, shape) -> int:
        """Compute max data rows that fit above the footer for *shape*.

        Uses each template data row's actual height (preserved when filling) plus
        the last template row's height as the per-new-row cost (new rows are cloned
        from the last row).  This avoids the old average-height approach that
        under-counted when the last row was significantly taller than average.
        """
        emu = self._EMU_PER_INCH
        default_row_h = emu // 3  # ~0.33" when a row has no explicit height

        try:
            avail_height = max(0, self._SVN_SAFE_BOTTOM_EMU - shape.top)

            if not shape.has_table:
                return max(1, avail_height // self._MIN_ROW_HEIGHT_EMU)

            table = shape.table
            all_rows = list(table.rows)
            if not all_rows:
                return max(1, avail_height // self._MIN_ROW_HEIGHT_EMU)

            header_h = all_rows[0].height or default_row_h
            data_rows = all_rows[1:] if len(all_rows) > 1 else all_rows

            # _apply_content_row_heights will bump every data row to at least
            # _MIN_DATA_ROW_HEIGHT_EMU after filling.  Use that same minimum here so
            # the pre-fill split threshold is consistent with the post-fill actual heights.
            effective_min = self._MIN_DATA_ROW_HEIGHT_EMU

            # Total height assuming template rows are bumped to at least effective_min.
            template_data_h = sum(
                max(r.height or default_row_h, effective_min) for r in data_rows
            )
            template_data_count = len(data_rows)

            # New rows cloned from last template row — also subject to effective_min.
            last_row_h = (data_rows[-1].height or default_row_h) if data_rows else default_row_h
            new_row_h = max(last_row_h, effective_min)

            avail_for_data = max(0, avail_height - header_h)
            remaining = max(0, avail_for_data - template_data_h)
            max_new = int(remaining / new_row_h) if new_row_h else 0
            result = template_data_count + max_new

            print(f'  - max_rows_for_shape: avail={avail_height/emu:.2f}" '
                  f'header={header_h/emu:.2f}" '
                  f'tmpl_data={template_data_h/emu:.2f}" ({template_data_count} rows) '
                  f'new_row={new_row_h/emu:.2f}" → {result} total data rows')
            return max(1, result)
        except Exception as e:
            fallback = max(1, self._SVN_SAFE_BOTTOM_EMU // self._MIN_ROW_HEIGHT_EMU)
            print(f'  - _max_rows_for_shape error ({e}); fallback={fallback}')
            return fallback

    def _normalize_table_col_widths(self, table, min_col_width_emu: int = 914400):
        """Ensure no column is narrower than *min_col_width_emu* (default 1 inch).

        When a column is below the minimum, it is expanded to the minimum and the
        surplus is subtracted proportionally from the wider columns, so the total
        table width stays the same.
        """
        try:
            cols = list(table.columns)
            if not cols:
                return
            widths = [c.width for c in cols]
            total_width = sum(widths)

            needs_expand = [w < min_col_width_emu for w in widths]
            if not any(needs_expand):
                return  # all columns already wide enough

            new_widths = list(widths)
            for i, below in enumerate(needs_expand):
                if below:
                    new_widths[i] = min_col_width_emu

            # Amount already locked to narrow columns
            locked_total = sum(w for i, w in enumerate(new_widths) if needs_expand[i])
            remaining = total_width - locked_total
            free_total = sum(w for i, w in enumerate(widths) if not needs_expand[i])

            if free_total > 0 and remaining > 0:
                scale = remaining / free_total
                for i, below in enumerate(needs_expand):
                    if not below:
                        new_widths[i] = max(min_col_width_emu, int(widths[i] * scale))

            # Fix rounding drift so total is exactly preserved
            diff = total_width - sum(new_widths)
            if diff != 0:
                # Apply diff to the widest free column
                widest = max((i for i, b in enumerate(needs_expand) if not b),
                             key=lambda i: new_widths[i], default=0)
                new_widths[widest] += diff

            for col, w in zip(cols, new_widths):
                col.width = w

            print(f'  - Normalized col widths: {[round(w / 914400, 2) for w in new_widths]}" each')
        except Exception as e:
            print(f'  - Error normalizing table col widths: {e}')

    def _fit_table_to_bounds(self, shape):
        """Constrain the table inside the slide's safe content area.

        After rows are added/cloned the table can grow taller than the area
        above the footer.  This method:
          1. Measures total row height after fill.
          2. If it fits → updates shape.height to the actual content height
             (removes dead space / avoids shape box drifting below table).
          3. If it overflows → scales all row heights proportionally down to
             fit within (safe_bottom – shape.top), with a per-row floor so
             text stays legible.
        """
        try:
            if not shape.has_table:
                return
            table = shape.table
            rows = list(table.rows)
            if not rows:
                return

            total_h = sum(r.height or 0 for r in rows)
            avail_h = max(0, self._SVN_SAFE_BOTTOM_EMU - shape.top)

            if total_h <= avail_h:
                shape.height = total_h   # sync shape box to actual content
                return

            # Need to compress — scale proportionally with a per-row floor
            min_h = self._MIN_ROW_HEIGHT_EMU
            scale = avail_h / total_h
            new_heights = [max(min_h, int((r.height or 0) * scale)) for r in rows]

            # Fix rounding drift: adjust the tallest row
            diff = avail_h - sum(new_heights)
            if diff != 0:
                idx = max(range(len(new_heights)), key=lambda i: new_heights[i])
                new_heights[idx] = max(min_h, new_heights[idx] + diff)

            for row, h in zip(rows, new_heights):
                row.height = h
            shape.height = sum(r.height for r in rows)
            print(f'  - Fit table to safe area: '
                  f'{total_h / 914400:.2f}" → {shape.height / 914400:.2f}" '
                  f'(safe bottom {self._SVN_SAFE_BOTTOM_EMU / 914400:.2f}")')
        except Exception as e:
            print(f'  - _fit_table_to_bounds error: {e}')

    def _find_shape(self, slide, target):
        """Find shape by index or name.

        When only ``shape_name`` is provided (no ``shape_index``), the search
        is performed recursively through GROUP shapes so that shapes nested
        inside groups can be targeted by name.
        """
        try:
            # Try by index first (faster, top-level only)
            if target.shape_index is not None and target.shape_index < len(slide.shapes):
                shape = slide.shapes[target.shape_index]
                print(f'  - Found shape by index {target.shape_index}: {shape.name}')
                return shape

            # Fallback to recursive name search (traverses group children)
            if target.shape_name:
                found = self._find_shape_recursive(slide.shapes, target.shape_name)
                if found:
                    print(f'  - Found shape by name (recursive): {target.shape_name}')
                    return found
        except Exception as e:
            print(f'  - Error finding shape: {e}')

        return None

    def _find_shape_recursive(self, shapes, name: str):
        """Recursively search for a shape by name, traversing GROUP children.

        Args:
            shapes: A shape collection (``slide.shapes`` or ``group_shape.shapes``).
            name: The exact shape name to search for.

        Returns:
            The matching shape, or ``None`` if not found.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in shapes:
                if shape.name == name:
                    return shape
                # Recurse into group shapes
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        found = self._find_shape_recursive(shape.shapes, name)
                        if found:
                            return found
                except AttributeError:
                    pass
        except Exception as e:
            print(f'  - _find_shape_recursive error: {e}')
        return None

    def _fill_slide_default(self, slide, slide_content: SlideContent):
        """Default filling strategy (original logic)"""

        # Update title
        self._update_slide_title(slide, slide_content.title)

        # Get available text containers
        placeholders = self._get_text_placeholders(slide)
        text_shapes = self._get_text_shapes(slide)

        # Fill content items
        placeholder_idx = 0
        for content_item in slide_content.content:
            if content_item.type == ContentType.TEXT:
                placeholder_idx = self._fill_text_content(
                    slide, content_item, placeholders, text_shapes, placeholder_idx,
                )

            elif content_item.type == ContentType.LIST:
                placeholder_idx = self._fill_list_content(
                    slide, content_item, placeholders, placeholder_idx,
                )

            elif content_item.type == ContentType.TABLE:
                self._add_table(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

            elif content_item.type == ContentType.IMAGE:
                self._add_image(slide, content_item, Inches(1.0), Inches(2.5), Inches(8.0))

    def _get_text_placeholders(self, slide):
        """Get all text placeholders in slide (excluding title)"""
        return [
            s for s in slide.shapes
            if s.is_placeholder and s.has_text_frame and s != slide.shapes.title
        ]

    def _get_text_shapes(self, slide):
        """Get all text shapes in slide (excluding placeholders and title)"""
        return [
            s for s in slide.shapes
            if s.has_text_frame and not s.is_placeholder and s != slide.shapes.title
        ]

    def _fill_text_content(self, slide, content_item, placeholders, text_shapes, placeholder_idx):
        """Fill text content into placeholders or text shapes"""
        if placeholder_idx < len(placeholders):
            placeholders[placeholder_idx].text_frame.text = content_item.data
            self._apply_text_font(placeholders[placeholder_idx].text_frame)
            print(f'  - Filled placeholder {placeholder_idx}: {content_item.data[:50]}...')
            return placeholder_idx + 1

        elif text_shapes:
            text_shapes[0].text_frame.text = content_item.data
            self._apply_text_font(text_shapes[0].text_frame)
            text_shapes.pop(0)
            return placeholder_idx

        else:
            # Add new textbox if no placeholders left
            self._add_content_at_position(slide, content_item, Inches(1.0), Inches(2.0), Inches(8.0))
            return placeholder_idx

    def _fill_list_content(self, slide, content_item, placeholders, placeholder_idx):
        """Fill list content into placeholder"""
        if placeholder_idx < len(placeholders):
            text_frame = placeholders[placeholder_idx].text_frame
            text_frame.clear()
            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
                p.level = 0
            self._apply_text_font(text_frame)
            print(f'  - Filled list with {len(content_item.data)} items')
            return placeholder_idx + 1
        else:
            # Add new list
            top = Inches(2.5 + (placeholder_idx * 0.5))
            self._add_content_at_position(slide, content_item, Inches(1.0), top, Inches(8.0))
            return placeholder_idx

    def _add_content_at_position(self, slide, content_item: ContentItem, left, top, width):
        """Add content item at specific position"""
        if content_item.type == ContentType.TEXT:
            height = Inches(0.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = content_item.data
            self._apply_text_font(textbox.text_frame)

        elif content_item.type == ContentType.LIST:
            height = Inches(0.3) * len(content_item.data)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            for i, item in enumerate(content_item.data):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = f'• {item}'
            self._apply_text_font(text_frame)

    def _add_image(self, slide, content_item: ContentItem, left, top, width):
        """Add image to slide"""
        try:
            image_data = content_item.data
            image_path = self._resolve_image_path(image_data)

            if Path(image_path).exists():
                slide.shapes.add_picture(str(image_path), left, top, width=width)
                self._add_image_caption(slide, image_data, left, top, width)
            else:
                print(f'Image not found: {image_path}')
                self._add_placeholder_text(slide, f'[Image: {image_path}]', left, top, width)

        except Exception as e:
            print(f'Error adding image: {e}')

    def _resolve_image_path(self, image_data) -> str:
        """Resolve image path from data"""
        image_path = image_data.get('path') if isinstance(image_data, dict) else image_data

        if image_path and not Path(str(image_path)).is_absolute():
            # Try relative to CWD first
            rel_path = Path(str(image_path))
            if not rel_path.exists():
                # Try outputs directory
                image_path = str(Path('outputs') / str(image_path))

        return str(image_path) if image_path else ''

    def _add_image_caption(self, slide, image_data, left, top, width):
        """Add caption below image if exists"""
        if isinstance(image_data, dict) and image_data.get('caption'):
            caption_top = top + Inches(3.0)
            caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.3))
            caption_box.text_frame.text = image_data['caption']
            for paragraph in caption_box.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Noto Sans JP'
                    run.font.size = Pt(10)
                    run.font.italic = True

    def _apply_text_font(self, text_frame, font_name: str = 'Noto Sans JP', font_size: int = 10):
        """Apply font and size to all text in a text frame"""
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to text: {e}')

    def _apply_line_spacing(self, text_frame, spacing_pt):
        """Apply line spacing (space after) to all paragraphs in a text frame
        
        Args:
            text_frame: The text frame to apply spacing to
            spacing_pt: Spacing in points (e.g., Pt(10) for 10pt)
        """
        try:
            for paragraph in text_frame.paragraphs:
                if paragraph.text.strip():  # skip empty paragraphs (from blank lines)
                    paragraph.space_after = spacing_pt
        except Exception as e:
            print(f'  - Error applying line spacing: {e}')

    def _estimate_max_chars(self, shape, font_pt: int = 8) -> int:
        """Estimate max characters that fit in a shape based on its dimensions.

        Uses CJK character width at the given font size to compute chars-per-line
        and shape height to compute how many lines fit.  Returns a floor of 8 so
        at least a short label is always written.
        """
        try:
            emu = self._EMU_PER_INCH
            w_inch = (shape.width or 0) / emu
            h_inch = (shape.height or 0) / emu
            cjk_w = font_pt / 72          # ~0.111" per CJK char at 8pt
            line_h = font_pt * 1.35 / 72  # ~0.150" line height with spacing
            chars_per_line = max(5, int(w_inch / cjk_w))
            lines = max(1, int(h_inch / line_h))
            return max(8, chars_per_line * lines)
        except Exception:
            return 65

    @staticmethod
    def _set_vertical_anchor_top(text_frame) -> None:
        """Anchor text to the top of the shape so it doesn't float to center."""
        try:
            from pptx.enum.text import MSO_ANCHOR
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass

    @staticmethod
    def _truncate_at_sentence(text: str, max_chars: int = 120) -> str:
        """Truncate text to fit within max_chars, preferring sentence boundaries.

        Tries Japanese sentence ends (。) then newlines then hard cut.
        Returns the original string unchanged when it is already short enough.
        """
        if len(text) <= max_chars:
            return text
        truncated = text[:max_chars]
        for sep in ['。', '\n', '．', '.']:
            idx = truncated.rfind(sep)
            if idx > max_chars // 2:
                return truncated[:idx + 1].rstrip()
        return truncated.rstrip() + '…'

    def _apply_auto_fit(self, text_frame):
        """Enable normAutofit so text shrinks to fit the shape instead of overflowing."""
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception as e:
            print(f'  - Could not apply auto-fit: {e}')

    def _apply_table_cell_font(self, cell, font_name: str = 'Noto Sans JP', font_size: int = 8):
        """Apply font and size to all text in a table cell"""
        try:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
        except Exception as e:
            print(f'  - Error applying font to table cell: {e}')

    def _add_placeholder_text(self, slide, text: str, left, top, width):
        """Add placeholder text box"""
        textbox = slide.shapes.add_textbox(left, top, width, Inches(1))
        textbox.text_frame.text = text
        self._apply_text_font(textbox.text_frame)

    def _add_table(self, slide, content_item: ContentItem, left, top, width):
        """Add table to slide"""
        try:
            table_data = content_item.data
            headers = table_data.get('headers', [])
            rows = table_data.get('rows', [])

            if not headers or not rows:
                print('Table has no headers or rows')
                return

            # Create table
            row_count = len(rows) + 1  # +1 for header
            col_count = len(headers)
            height = Inches(0.3) * row_count

            table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
            table = table_shape.table

            # Fill headers
            self._fill_table_headers(table, headers)

            # Fill data rows
            self._fill_table_rows(table, rows)

        except Exception as e:
            print(f'Error adding table: {e}')

    def _fill_table_headers(self, table, headers: List[str]):
        """Fill table header row"""
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            # Apply font and make header bold
            self._apply_table_cell_font(cell)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    def _fill_table_rows(self, table, rows: List[List[str]]):
        """Fill table data rows"""
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)  # +1 to skip header
                cell.text = str(cell_data)
                self._apply_table_cell_font(cell)

    # ==================================================================
    # Role-based rendering (new pipeline: profile.md -> ProjectProfile -> PPTX)
    # ==================================================================

    def render_from_profile(
        self,
        profile: ProjectProfile,
        template: str,
        output_name: str,
        output_dir: Optional[str] = None,
        extra_slides: Optional[List[dict]] = None,
    ) -> Path:
        """Fill template using a ProjectProfile + role-based slide configs.

        For each configured slide:
          - Resolve each role's data via source_path
          - Discover shapes via ShapeDiscovery
          - Fill shapes (text / table / image / compound pairs)
          - Clear excess shapes
        """
        print(f'Rendering from profile: {output_name}')
        self._output_dir = self._resolve_output_dir(output_dir)
        self._detect_template_config(template)
        prs = self._load_template(template)
        # Compute safe content area from actual slide height so tables never reach the footer.
        # Footer band occupies roughly the bottom 13% of the slide; 0.87 leaves ~0.36" margin
        # above a footer at 93% (observed: 5.254" on a 5.625" slide → 93.4%).
        if prs.slide_height:
            self._SVN_SAFE_BOTTOM_EMU = int(prs.slide_height * 0.87)
            print(f'  Safe content bottom: {self._SVN_SAFE_BOTTOM_EMU / self._EMU_PER_INCH:.3f}" '
                  f'(slide is {prs.slide_height / self._EMU_PER_INCH:.3f}" tall)')
        if prs.slide_width:
            self._slide_width_emu = prs.slide_width
        last_positions = self._fill_slides_from_profile(prs, profile)
        if extra_slides:
            self._render_extra_slides(prs, extra_slides, last_positions)
        return self._save_presentation(prs, output_name, output_dir)

    def _fill_slides_from_profile(self, prs: Presentation, profile: ProjectProfile) -> Dict[int, int]:
        """Fill configured slides; return map {template_slide_num: final_0based_index_of_last_variant}."""
        last_positions: Dict[int, int] = {}
        if not self.template_config:
            print('No template config — cannot render from profile')
            return last_positions

        configs = self.template_config.slide_role_configs
        sorted_nums = sorted(configs.keys())
        slide_offset = 0
        discovery = ShapeDiscovery()

        for slide_num in sorted_nums:
            if not self.template_config.should_fill_slide(slide_num):
                print(f'Slide {slide_num} is protected — skipping')
                continue
            adj_idx = slide_num - 1 + slide_offset
            if adj_idx >= len(prs.slides):
                print(f'Slide {slide_num} not found in template (have {len(prs.slides)})')
                continue
            slide = prs.slides[adj_idx]
            role_config = configs[slide_num]
            print(f'Filling slide {slide_num} (template index {adj_idx}): '
                  f'{len(role_config.roles)} role(s)')
            extra = self._fill_slide_roles(prs, slide, role_config, profile, adj_idx, discovery)
            # Last position of this section's variant chain (after overflow inserts)
            last_positions[slide_num] = adj_idx + extra
            slide_offset += extra
        return last_positions

    def _fill_slide_roles(
        self, prs: Presentation, slide, role_config: SlideRoleConfig,
        profile: ProjectProfile, slide_index: int, discovery: ShapeDiscovery,
    ) -> int:
        """Fill all roles on a slide. Returns extra slides inserted (table overflow)."""
        extra_total = 0
        for role in role_config.roles:
            data = self._resolve_source_path(profile, role.source_path)
            shapes = discovery.discover(slide, role)
            if role.kind == ContentKind.TABLE:
                extra_total += self._fill_table_role(prs, slide, slide_index + extra_total,
                                                    shapes, data, role)
            elif role.kind == ContentKind.IMAGE:
                self._fill_image_role(slide, shapes, data, role)
            else:  # TEXT (simple or compound)
                if role.sub_roles:
                    self._fill_compound_role(shapes, data, role)
                else:
                    self._fill_text_role(shapes, data, role)
        return extra_total

    # ------------------------------------------------------------------
    # Source-path resolution: dot notation + optional [start:end] slice
    # ------------------------------------------------------------------

    def _resolve_source_path(self, profile: ProjectProfile, path: str):
        """Resolve `a.b.c` or `a[0:2]` or `a.b[2:4]` into a value on the profile.

        Returns None if any segment is missing.
        """
        if not path:
            return None
        # Split off trailing slice like benefits[0:2]
        slice_spec = None
        if path.endswith(']') and '[' in path:
            base, slc = path.rsplit('[', 1)
            slice_spec = slc[:-1]  # strip ']'
            path = base
        cur = profile
        for part in path.split('.'):
            if cur is None:
                return None
            cur = getattr(cur, part, None)
        if slice_spec is not None and isinstance(cur, list):
            try:
                a, b = slice_spec.split(':')
                start = int(a) if a else None
                end = int(b) if b else None
                return cur[start:end]
            except Exception:
                return cur
        return cur

    # ------------------------------------------------------------------
    # Role fillers
    # ------------------------------------------------------------------

    def _fill_text_role(self, shapes: list, data, role: ShapeRole):
        """Fill text into a list of shapes. data: str | list[str]."""
        items = self._coerce_text_items(data)
        for i, shape in enumerate(shapes):
            if shape is None:
                continue
            value = items[i] if i < len(items) else ''
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True
                shape.text_frame.text = value
                self._apply_text_font(shape.text_frame)
                self._apply_auto_fit(shape.text_frame)
                if role.name == 'current_issues' or role.name == 'objectives':
                    self._apply_line_spacing(shape.text_frame, Pt(10))

    def _fill_compound_role(self, pairs: list, data, role: ShapeRole):
        """Fill compound (title, body) pairs. data: list[obj] with .title/.body or dict."""
        items = data if isinstance(data, list) else []
        for i, pair in enumerate(pairs):
            if i >= len(items):
                # Clear unused pair shapes
                if isinstance(pair, tuple):
                    for shp in pair:
                        ShapeDiscovery.clear_text(shp)
                continue
            item = items[i]
            title = getattr(item, 'title', None) or (item.get('title') if isinstance(item, dict) else '')
            body = getattr(item, 'body', None) or getattr(item, 'content', None) \
                or (item.get('body') if isinstance(item, dict) else '') \
                or (item.get('content') if isinstance(item, dict) else '')
            if not isinstance(pair, tuple):
                continue
            title_shape, body_shape = pair
            if title_shape is not None and title_shape.has_text_frame:
                title_shape.text_frame.word_wrap = True
                title_shape.text_frame.text = title or ''
                self._apply_text_font(title_shape.text_frame)
                self._set_vertical_anchor_top(title_shape.text_frame)
            if body_shape is not None and body_shape.has_text_frame:
                max_chars = self._estimate_max_chars(body_shape)
                body_text = self._truncate_at_sentence(body or '', max_chars=max_chars)
                body_shape.text_frame.word_wrap = True
                body_shape.text_frame.text = body_text
                self._apply_text_font(body_shape.text_frame)
                self._set_vertical_anchor_top(body_shape.text_frame)

    def _fill_table_role(
        self, prs: Presentation, slide, slide_index: int,
        shapes: list, data, role: ShapeRole,
    ) -> int:
        """Fill a table role. Reuses overflow logic by converting to markdown."""
        if not shapes or shapes[0] is None:
            print(f"  - Table shape not found for role '{role.name}'")
            return 0
        shape = shapes[0]
        rows = data if isinstance(data, list) else []
        if not rows:
            return 0
        md_table = self._rows_to_md_table(rows, fill_cols=role.fill_cols)
        target = _LegacyTarget(content_key=role.name, fill_cols=role.fill_cols)
        # Use shape-position-aware row limit so tables never overflow into footer
        max_rows = self._max_rows_for_shape(shape)
        chunks = self._split_table_data(md_table, max_rows)
        if len(chunks) <= 1:
            self._fill_table(shape, md_table, target)
            return 0

        print(f'  - Role "{role.name}": table split ({len(chunks)} slides, max {max_rows} rows/slide)')
        self._fill_table(shape, chunks[0], target)
        extra = 0
        current_index = slide_index
        for i, chunk in enumerate(chunks[1:], 1):
            new_slide = self._duplicate_slide(prs, current_index)
            extra += 1
            current_index += 1
            # Re-discover same shape on duplicate, then fill
            new_shape = ShapeDiscovery().discover(new_slide, role)
            if new_shape and new_shape[0] is not None:
                self._fill_table(new_shape[0], chunk, target)
                print(f'  - Filled continuation {i}/{len(chunks)-1}')
        return extra

    def _fill_image_role(self, slide, shapes: list, data, role: ShapeRole):
        """Fill an image role. data: str (image path or data URI)."""
        if not shapes or shapes[0] is None:
            print(f"  - Image shape not found for role '{role.name}'")
            return
        if not data:
            return
        if getattr(role, 'fit_to_slide', False):
            self._place_image_fit_slide(slide, data, role.name)
        else:
            target = _LegacyTarget(content_key=role.name)
            self._replace_shape_with_image(slide, shapes[0], data, target)

    def _place_image_fit_slide(self, slide, image_path: str, role_name: str):
        """Scale image to fill slide content area (below title), preserving aspect ratio.

        Uses Pillow to read native image dimensions.  Falls back to placeholder-based
        placement when the file cannot be opened.
        """
        image_path = image_path.strip()

        # Resolve relative path
        if not Path(image_path).is_absolute():
            for candidate in [Path(image_path),
                               Path('outputs') / image_path,
                               (self._output_dir / image_path) if self._output_dir else None]:
                if candidate and candidate.exists():
                    image_path = str(candidate)
                    break

        if not Path(image_path).exists():
            print(f'  - Image not found for fit-to-slide: {image_path}')
            return

        try:
            from PIL import Image as _PILImage
            with _PILImage.open(image_path) as img:
                img_w, img_h = img.size
        except Exception as e:
            print(f'  - Could not read image dimensions ({e}); skipping fit-to-slide')
            return

        # Use actual slide dimensions recorded during render_from_profile.
        # Fall back to common 10" width if not yet set.
        slide_w = self._slide_width_emu or 9144000
        # Use the dynamically computed safe-bottom (above footer) as the effective
        # lower boundary so the image never overlaps the footer band.
        safe_bottom = self._SVN_SAFE_BOTTOM_EMU
        margin_l = int(0.5 * self._EMU_PER_INCH)   # 0.5"
        margin_r = int(0.5 * self._EMU_PER_INCH)   # 0.5"
        margin_t = int(1.1 * self._EMU_PER_INCH)   # 1.1" (below title bar)

        avail_w = slide_w - margin_l - margin_r
        avail_h = max(0, safe_bottom - margin_t)

        scale = min(avail_w / img_w, avail_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)

        # Center within the content area
        left = margin_l + (avail_w - new_w) // 2
        top = margin_t + (avail_h - new_h) // 2

        try:
            slide.shapes.add_picture(image_path, left, top, width=new_w, height=new_h)
            print(f"  Placed '{role_name}' fit-to-slide: "
                  f"{new_w / 914400:.2f}\" × {new_h / 914400:.2f}\" "
                  f"at ({left / 914400:.2f}\", {top / 914400:.2f}\")")
        except Exception as e:
            print(f'  - Error placing fit-to-slide image: {e}')

    # ------------------------------------------------------------------
    # Helpers for role-based fill
    # ------------------------------------------------------------------

    @staticmethod
    def _coerce_text_items(data) -> list:
        """Normalize text-role data into a list of strings."""
        if data is None:
            return []
        if isinstance(data, str):
            return [data]
        if isinstance(data, list):
            return [str(item) if item is not None else '' for item in data]
        return [str(data)]

    @staticmethod
    def _rows_to_md_table(rows: list, fill_cols=None) -> str:
        """Convert list[dict | dataclass] to markdown pipe-table string.

        - When fill_cols is None: emit standard header + data rows.
        - When fill_cols is set: emit row-only format (no header/separator), and
          emit only the *last* len(fill_cols) fields per row (template provides
          fixed columns for the rest, e.g. assumption labels).
        """
        if not rows:
            return ''
        # Normalize rows to list[dict]
        import dataclasses as _dc
        rows = [r if isinstance(r, dict) else _dc.asdict(r) if _dc.is_dataclass(r) else dict(vars(r))
                for r in rows]
        headers = list(rows[0].keys())

        def esc(v):
            return str(v).replace('|', '&#124;').replace('\n', ' ')

        if fill_cols:
            n = len(fill_cols)
            target_keys = headers[-n:] if n <= len(headers) else headers
            lines = []
            for r in rows:
                cells = [esc(r.get(k, '')) for k in target_keys]
                lines.append('| ' + ' | '.join(cells) + ' |')
            return '\n'.join(lines)

        lines = [
            '| ' + ' | '.join(headers) + ' |',
            '|' + '|'.join(['---'] * len(headers)) + '|',
        ]
        for r in rows:
            lines.append('| ' + ' | '.join(esc(r.get(h, '')) for h in headers) + ' |')
        return '\n'.join(lines)


    # ==================================================================
    # Extra-slide rendering — AI-generated, generate-slide-style layouts
    # ==================================================================

    def _render_extra_slides(
        self, prs: Presentation, extras: List[dict], last_positions: Dict[int, int],
    ) -> None:
        """Insert AI-generated slides rendered with generate-slide-style layouts.

        Each entry of *extras* is a dict with keys:
          - layout:          str  (numbered_points|card_grid|comparison_2|... ; default bullets)
          - title:           str  (slide title)
          - <layout content> (e.g. points / items / columns / steps / message)
          - anchor_section:  str | None (profile section name; resolves to anchor slide)
          - anchor_slide:    int | None (explicit slide number; overrides anchor_section)

        New slides are created from a CONTENT slide's layout (so they inherit the
        deck's white background + master footer/page-number/logo) and inserted
        right after the LAST variant of the anchor slide (accounting for overflow
        inserts). Multiple extras sharing an anchor keep their given order.
        """
        from lib.templates.svn import SECTION_TO_LAST_SLIDE

        w_emu, h_emu = prs.slide_width, prs.slide_height

        # Prefer a chrome reference slide — a configured CONTENT slide that
        # has slide-level chrome shapes (logo / footer in the bottom strip).
        # Cloning that slide makes the chrome real slide-level shapes (editable)
        # rather than master-only chrome (which appears "captured" / locked).
        chrome_ref_idx = self._chrome_canvas_idx(prs, last_positions)
        canvas_layout = None
        if chrome_ref_idx is None:
            # Fallback: master-chrome only (legacy behavior).
            canvas_layout = self._content_canvas_layout(prs, last_positions)
            if canvas_layout is None:
                print('No content-slide layout available for extra slides')
                return

        # Group by anchor (preserve given order within a group)
        groups: Dict[int, List[dict]] = {}
        for e in extras:
            anchor = self._resolve_extra_anchor(e, last_positions, SECTION_TO_LAST_SLIDE, prs)
            groups.setdefault(anchor, []).append(e)

        # Insertions shift later positions; process in ascending anchor order
        # and track cumulative offset.
        offset = 0
        for anchor in sorted(groups.keys()):
            for i, extra in enumerate(groups[anchor]):
                target = anchor + offset + 1  # 0-indexed insertion point
                if chrome_ref_idx is not None:
                    # Source index for chrome_ref shifts right with every prior
                    # insert that landed at or before it.
                    src_idx = chrome_ref_idx
                    new_slide = self._duplicate_slide_at(prs, src_idx, target)
                    self._strip_content_keep_chrome(new_slide, h_emu)
                    self._fill_extra_slide(new_slide, extra, w_emu, h_emu, clear=False)
                    if src_idx >= target:
                        chrome_ref_idx += 1  # our insert shifted the ref slide
                else:
                    new_slide = self._add_slide_at(prs, canvas_layout, target)
                    self._fill_extra_slide(new_slide, extra, w_emu, h_emu)
                offset += 1
                print(f'  - Inserted extra slide "{extra.get("title", "")[:40]}"'
                      f' [{extra.get("layout", "bullets")}] at index {target}')

    def _content_canvas_layout(self, prs: Presentation, last_positions: Dict[int, int]):
        """Pick a CLEAN content-slide layout (light bg + standard small footer logo).

        Some configured slides (e.g. slide 4 in the SVN deck) use a unique master
        with a large decorative side panel and an oversized logo. Using that
        layout for AI-generated extra slides makes them visually inconsistent
        with the rest of the deck. We skip those and prefer the layout of a
        configured content slide whose master only contains small chrome
        (footer text, page number, small logo).
        """
        if not last_positions:
            return prs.slides[0].slide_layout if len(prs.slides) else None

        # Iterate configured slides in template-order; pick the first clean one.
        slide_w = prs.slide_width or 0
        slide_h = prs.slide_height or 0
        slide_area = slide_w * slide_h
        for sn in sorted(last_positions):
            idx = last_positions[sn]
            if not (0 <= idx < len(prs.slides)):
                continue
            layout = prs.slides[idx].slide_layout
            if self._is_clean_master(layout.slide_master, slide_area):
                return layout

        # Fallback: earliest configured slide's layout (legacy behavior)
        idx = last_positions[min(last_positions)]
        if 0 <= idx < len(prs.slides):
            return prs.slides[idx].slide_layout
        return prs.slides[0].slide_layout if len(prs.slides) else None

    @staticmethod
    def _is_clean_master(master, slide_area: int) -> bool:
        """A master is 'clean' if it has no large solid-filled decorative shapes.

        Large = >= 25% of the slide area. This filters out section-divider style
        masters (e.g. SVN slide 4 with a pink right-half panel) while keeping
        standard content-slide masters that only carry footer / logo chrome.
        """
        if slide_area <= 0:
            return True
        threshold = slide_area * 0.25
        try:
            for sh in master.shapes:
                if not (sh.width and sh.height):
                    continue
                if sh.width * sh.height < threshold:
                    continue
                # Large shape — clean only if it's not a solid colored fill
                try:
                    fill = sh.fill
                    if fill.type == 1:  # MSO_FILL.SOLID
                        return False
                except Exception:
                    pass
        except Exception:
            return True
        return True

    def _add_slide_at(self, prs: Presentation, layout, target_index: int):
        """Add a fresh slide from *layout* and move it to target_index (0-based)."""
        new_slide = prs.slides.add_slide(layout)
        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            sldIdLst = prs.part._element.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                new_sldId = list(sldIdLst)[-1]  # just appended → at end
                sldIdLst.remove(new_sldId)
                remaining = list(sldIdLst)
                if target_index >= len(remaining):
                    sldIdLst.append(new_sldId)
                else:
                    remaining[target_index].addprevious(new_sldId)
        except Exception as e:
            print(f'  - Error positioning extra slide: {e}')
        return prs.slides[target_index]

    def _resolve_extra_anchor(
        self, extra: dict, last_positions: Dict[int, int],
        section_map: Dict[str, int], prs: Presentation,
    ) -> int:
        """Return 0-indexed slide position to insert AFTER."""
        explicit = extra.get('anchor_slide')
        if isinstance(explicit, int) and explicit > 0:
            return last_positions.get(explicit, explicit - 1)
        section = extra.get('anchor_section')
        if section and section in section_map:
            slide_num = section_map[section]
            return last_positions.get(slide_num, slide_num - 1)
        # Default: append before the extra-slide template (i.e. after last filled slide)
        if last_positions:
            return max(last_positions.values())
        return len(prs.slides) - 2  # before the template slide itself

    def _fill_extra_slide(self, slide, extra: dict, w_emu: int, h_emu: int,
                          clear: bool = True) -> None:
        """Render a generate-slide-style layout onto the cloned canvas.

        When *clear* is True the slide is cleared first (use for slides created
        from a layout). When False, existing shapes are kept — used after the
        slide was duplicated from an anchor and its content shapes stripped, so
        only the chrome (logo / footer) remains and the layout is drawn on top.

        See lib/extra_slide_layouts.py for supported layouts.
        """
        from lib.extra_slide_layouts import render_extra_layout
        render_extra_layout(slide, extra, w_emu, h_emu, clear=clear)

    # ------------------------------------------------------------------
    # Chrome-aware slide cloning helpers (used for extra AI-generated slides)
    # ------------------------------------------------------------------

    def _chrome_canvas_idx(self, prs: Presentation,
                           last_positions: Dict[int, int]) -> Optional[int]:
        """Find a configured content slide that has slide-level chrome shapes.

        Chrome shapes are small elements parked in the bottom strip of the
        slide (logo, footer text). Cloning such a slide preserves the chrome as
        real slide-level shapes — editable in PowerPoint, rather than appearing
        as locked master-only chrome.

        Returns the slide index, or None when no suitable slide is found.
        """
        if not last_positions:
            return None
        slide_w = prs.slide_width or 0
        slide_h = prs.slide_height or 0
        if slide_h <= 0:
            return None
        slide_area = slide_w * slide_h
        threshold_top = int(slide_h * 0.88)  # shapes anchored in bottom ~12%
        # Iterate configured slides in template order, prefer clean masters.
        for sn in sorted(last_positions):
            idx = last_positions[sn]
            if not (0 <= idx < len(prs.slides)):
                continue
            slide = prs.slides[idx]
            if not self._is_clean_master(slide.slide_layout.slide_master, slide_area):
                continue
            for sh in slide.shapes:
                if sh.top is not None and sh.top >= threshold_top:
                    return idx
        return None

    def _duplicate_slide_at(self, prs: Presentation, source_idx: int,
                            target_idx: int):
        """Duplicate slide at *source_idx* and insert the copy at *target_idx*.

        Differs from :meth:`_duplicate_slide` (which always inserts immediately
        after the source) in that it places the copy at an arbitrary position.
        Also re-binds any relationship references (``r:embed`` / ``r:link``
        / ``r:id``) inside the cloned shapes so embedded images survive the
        cross-slide copy — otherwise the new slide's blipFills would resolve
        their rIds against the wrong relationship table.
        """
        source_slide = prs.slides[source_idx]

        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        src_tree = source_slide.shapes._spTree
        new_tree = new_slide.shapes._spTree
        for child in list(new_tree):
            new_tree.remove(child)
        for child in src_tree:
            new_tree.append(copy.deepcopy(child))

        self._remap_shape_rels(source_slide.part, new_slide.part, new_tree)

        pml_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        try:
            sldIdLst = prs.part._element.find(f'{{{pml_ns}}}sldIdLst')
            if sldIdLst is not None:
                new_sldId = list(sldIdLst)[-1]  # just appended at end
                sldIdLst.remove(new_sldId)
                remaining = list(sldIdLst)
                if target_idx >= len(remaining):
                    sldIdLst.append(new_sldId)
                else:
                    remaining[target_idx].addprevious(new_sldId)
        except Exception as e:
            print(f'  - Error positioning extra slide: {e}')
        return prs.slides[target_idx]

    @staticmethod
    def _remap_shape_rels(source_part, target_part, new_spTree) -> None:
        """Re-bind ``r:embed`` / ``r:link`` / ``r:id`` references in *new_spTree*
        to relationships on *target_part*.

        Cross-slide deep-copies preserve the textual rId values (e.g. ``rId1``)
        but those IDs are scoped to the originating slide's relationship table.
        Without remapping, an image blipFill copied from another slide ends up
        pointing at whatever the destination slide's ``rId1`` happens to be
        (typically the slideLayout), causing the image to vanish.
        """
        r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        src_rels = source_part.rels
        rid_map: Dict[str, str] = {}
        for el in new_spTree.iter():
            for attr_name in list(el.attrib):
                if not attr_name.startswith(f'{{{r_ns}}}'):
                    continue
                local = attr_name.split('}', 1)[1]
                if local not in ('embed', 'link', 'id'):
                    continue
                old_rid = el.attrib[attr_name]
                if not old_rid:
                    continue
                if old_rid in rid_map:
                    el.attrib[attr_name] = rid_map[old_rid]
                    continue
                try:
                    rel = src_rels[old_rid]
                except KeyError:
                    continue
                try:
                    if rel.is_external:
                        new_rid = target_part.relate_to(
                            rel.target_ref, rel.reltype, is_external=True)
                    else:
                        new_rid = target_part.relate_to(
                            rel.target_part, rel.reltype)
                except Exception as exc:
                    print(f'  - Warn: could not rebind {old_rid} ({rel.reltype}): {exc}')
                    continue
                rid_map[old_rid] = new_rid
                el.attrib[attr_name] = new_rid

    @staticmethod
    def _strip_content_keep_chrome(slide, slide_h_emu: int) -> None:
        """Remove every shape whose top edge sits above the bottom chrome strip.

        Anything anchored in the bottom ~12% of the slide is treated as chrome
        (logo / footer text) and retained so the layout renderer can draw the
        new content on top while the chrome remains intact and editable.
        """
        if slide_h_emu <= 0:
            return
        threshold = int(slide_h_emu * 0.88)
        for sh in list(slide.shapes):
            if sh.top is None or sh.top < threshold:
                sh._element.getparent().remove(sh._element)


class _LegacyTarget:
    """Lightweight stand-in for legacy ShapeTarget used by _fill_table / _replace_shape_with_image."""

    def __init__(self, content_key: str, fill_cols=None):
        self.content_key = content_key
        self.fill_cols = fill_cols
        self.shape_index = None
        self.shape_name = None
