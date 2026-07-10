"""Render generate-slide-style layouts onto a python-pptx slide (SVN deck canvas).

Ports a subset of the Sun* slide layout library
(claude/skills/generate-slide/references/layouts.md) to python-pptx so that
AI-authored "extra" slides match the polished generate-slide look while living
inside the SVN proposal deck.

Coordinates in this module are authored for the generate-slide LAYOUT_WIDE
canvas (13.333 x 7.5 in) and scaled automatically to the target slide size
(SVN template is 10 x 5.625 in). Brand constants mirror
generate-slide/assets/helpers.js.

Supported layouts (entry['layout']):
  - bullets        (default) : {bullets: [str, ...]}
  - numbered_points          : {points: [{header, body}, ...]}  (3-5)
  - card_grid                : {items: [{title, body, icon?}, ...]} (up to 4)
  - comparison_2             : {columns: [{title, items: [str]}, {title, items: [str]}]}
  - process_flow             : {steps: [str, ...]} (3-5)
  - hero                     : {message: str}
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# Brand constants (mirror generate-slide/assets/helpers.js)
SUN_RED = RGBColor(0xFF, 0x22, 0x00)
SUN_DARK_RED = RGBColor(0xAD, 0x0C, 0x00)
SUN_GOLD = RGBColor(0xB6, 0x92, 0x56)
LIGHT_GREY = RGBColor(0xF7, 0xF7, 0xF7)
BORDER_GREY = RGBColor(0xDD, 0xDD, 0xDD)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Noto Sans JP'  # single font keeps CJK/Vietnamese safe

GEN_W, GEN_H = 13.333, 7.5  # generate-slide authoring canvas (inches)
_EMU_IN = 914400


class _Canvas:
    """Coordinate/font scaler + shape helpers bound to one slide."""

    def __init__(self, slide, w_emu: int, h_emu: int):
        self.slide = slide
        self.sx = (w_emu / _EMU_IN) / GEN_W
        self.sy = (h_emu / _EMU_IN) / GEN_H
        self.fs = self.sy  # font scale follows vertical scale

    def _x(self, inch: float) -> Emu:
        return Emu(int(inch * self.sx * _EMU_IN))

    def _y(self, inch: float) -> Emu:
        return Emu(int(inch * self.sy * _EMU_IN))

    def rect(self, x, y, w, h, color: RGBColor):
        sp = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, self._x(x), self._y(y), self._x(w), self._y(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def oval(self, x, y, w, h, color: RGBColor):
        sp = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL, self._x(x), self._y(y), self._x(w), self._y(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def text(self, x, y, w, h, value, size, color, *, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
        tb = self.slide.shapes.add_textbox(self._x(x), self._y(y), self._x(w), self._y(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        lines = value if isinstance(value, list) else [value]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = str(line)
            run.font.name = FONT
            run.font.size = Pt(max(7, size * self.fs))
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
        return tb


def render_extra_layout(slide, entry: dict, w_emu: int, h_emu: int,
                        clear: bool = True) -> None:
    """Render the chosen layout for one extra slide.

    When *clear* is True (default), the slide's own shapes are removed first —
    suitable for a fresh slide created from a content-slide layout where chrome
    comes from the master.

    When *clear* is False, existing shapes are kept (caller is expected to have
    already left only the chrome shapes on the slide, e.g. via duplicating an
    anchor slide then stripping the content shapes). This allows the chrome
    (logo / footer) to remain as editable slide-level shapes rather than only
    living on the master.
    """
    if clear:
        _clear_shapes(slide)
    cv = _Canvas(slide, w_emu, h_emu)

    title = (entry.get('title') or '').strip()
    layout = (entry.get('layout') or 'bullets').strip().lower()

    if title:
        cv.text(0.7, 0.45, 12.0, 0.7, title, 26, TEXT_DARK, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
        cv.rect(0.7, 1.15, 1.6, 0.06, SUN_RED)

    renderer = _LAYOUTS.get(layout, _bullets)
    renderer(cv, entry)


def _clear_shapes(slide) -> None:
    """Remove the canvas slide's own shapes (master chrome is unaffected)."""
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


# ---------------------------------------------------------------------------
# Layout renderers — coordinates from layouts.md (13.333 x 7.5 canvas)
# ---------------------------------------------------------------------------

def _bullets(cv: _Canvas, entry: dict) -> None:
    bullets = entry.get('bullets') or entry.get('content', {}).get('bullets') or []
    if isinstance(bullets, str):
        bullets = [b.strip() for b in bullets.splitlines() if b.strip()]
    lines = [f'■  {b}' for b in bullets]
    cv.text(0.9, 1.7, 11.5, 5.0, lines, 15, TEXT_DARK)


def _numbered_points(cv: _Canvas, entry: dict) -> None:
    points = _get(entry, 'points')
    n = max(1, len(points))
    step = 1.4 if n <= 4 else 0.95
    for i, p in enumerate(points[:5]):
        y = 1.55 + i * step
        cv.oval(0.7, y, 0.6, 0.6, SUN_RED)
        cv.text(0.7, y, 0.6, 0.6, str(i + 1), 20, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cv.text(1.5, y - 0.05, 11.0, 0.45, _f(p, 'header'), 16, TEXT_DARK,
                bold=True, anchor=MSO_ANCHOR.MIDDLE)
        body = _f(p, 'body')
        if body:
            cv.text(1.5, y + 0.45, 11.0, step - 0.5, body, 12, TEXT_GREY)


def _card_grid(cv: _Canvas, entry: dict) -> None:
    items = _get(entry, 'items')[:4]
    for i, item in enumerate(items):
        col, row = i % 2, i // 2
        x = 0.6 + col * 6.2
        y = 1.5 + row * 2.35
        cv.rect(x, y, 6.0, 2.1, LIGHT_GREY)
        cv.rect(x, y, 0.1, 2.1, SUN_RED)
        icon = _f(item, 'icon')
        tx = x + (0.95 if icon else 0.3)
        if icon:
            cv.text(x + 0.25, y + 0.15, 0.7, 0.7, icon, 30, TEXT_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        cv.text(tx, y + 0.2, 4.8, 0.45, _f(item, 'title'), 15, TEXT_DARK,
                bold=True, anchor=MSO_ANCHOR.MIDDLE)
        cv.text(x + 0.3, y + 0.95, 5.4, 1.0, _f(item, 'body'), 11, TEXT_GREY)


def _comparison_2(cv: _Canvas, entry: dict) -> None:
    cols = _get(entry, 'columns')[:2]
    accents = [SUN_GOLD, SUN_RED]
    bgs = [RGBColor(0xFA, 0xF5, 0xEC), RGBColor(0xFF, 0xEE, 0xEC)]
    col_w, start_x, start_y, col_h = 6.0, 0.6, 1.55, 4.9
    for i, col in enumerate(cols):
        x = start_x + i * (col_w + 0.3)
        cv.rect(x, start_y, col_w, col_h, bgs[i])
        cv.rect(x, start_y, 0.1, col_h, accents[i])
        cv.text(x + 0.35, start_y + 0.2, col_w - 0.6, 0.55, _f(col, 'title'),
                19, accents[i], bold=True, anchor=MSO_ANCHOR.MIDDLE)
        items = col.get('items') or []
        lines = [f'■  {it}' for it in items]
        cv.text(x + 0.35, start_y + 1.0, col_w - 0.6, col_h - 1.2, lines, 12, TEXT_DARK)


def _process_flow(cv: _Canvas, entry: dict) -> None:
    steps = _get(entry, 'steps')[:5]
    if isinstance(steps, str):
        steps = [s.strip() for s in steps.splitlines() if s.strip()]
    n = max(1, len(steps))
    gap = 0.2
    total = GEN_W - 1.2
    step_w = (total - gap * (n - 1)) / n
    start_x = 0.6
    y = 3.0
    for i, s in enumerate(steps):
        x = start_x + i * (step_w + gap)
        cv.rect(x, y, step_w, 1.5, LIGHT_GREY)
        cv.rect(x, y, 0.1, 1.5, SUN_RED)
        cv.text(x + 0.25, y + 0.15, step_w - 0.4, 0.5, str(i + 1), 18, SUN_RED,
                bold=True)
        cv.text(x + 0.25, y + 0.65, step_w - 0.45, 0.75, _f2(s), 11, TEXT_DARK)


def _hero(cv: _Canvas, entry: dict) -> None:
    msg = entry.get('message') or entry.get('content', {}).get('message') or ''
    cv.rect(0.7, 3.0, 12.0, 1.6, SUN_RED)
    cv.text(0.9, 3.0, 11.6, 1.6, msg, 22, WHITE, bold=True, italic=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


_LAYOUTS = {
    'bullets': _bullets,
    'numbered_points': _numbered_points,
    'card_grid': _card_grid,
    'comparison_2': _comparison_2,
    'process_flow': _process_flow,
    'hero': _hero,
}


# ---------------------------------------------------------------------------
# Small helpers — tolerate both top-level keys and a nested `content` dict
# ---------------------------------------------------------------------------

def _get(entry: dict, key: str) -> list:
    return entry.get(key) or (entry.get('content') or {}).get(key) or []


def _f(obj, key: str) -> str:
    if isinstance(obj, dict):
        return str(obj.get(key) or '')
    return str(obj) if key in ('title', 'header') else ''


def _f2(obj) -> str:
    if isinstance(obj, dict):
        return str(obj.get('text') or obj.get('title') or obj.get('header') or '')
    return str(obj)
