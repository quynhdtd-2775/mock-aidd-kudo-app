"""Shape discovery engine — locate template shapes for a given ShapeRole.

Strategies:
  BY_INDEX           — direct slide.shapes[i]; or list of indices for N-cardinality
                       or pair_indices=[(title_idx, body_idx), ...] for compound roles
  BY_NAME            — recursive search by exact shape name(s); supports pair_names
                       for compound roles
  BY_GROUP_CHILDREN  — find a group shape by name, return children sorted by position
  BY_POSITION        — filter by shape type, sort by axis ('top' or 'left'), take N

Returns lists of pptx Shape objects in the order they should be filled.
Compound roles return list[(title_shape, body_shape)] tuples.
"""
from __future__ import annotations

from typing import Any, List, Optional, Tuple

try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    MSO_SHAPE_TYPE = None

from lib.templates.base import ContentKind, ShapeFindStrategy, ShapeRole


class ShapeDiscovery:
    """Locate shapes on a slide that fulfill a ShapeRole."""

    def discover(self, slide, role: ShapeRole) -> List[Any]:
        """Return ordered shape list for this role.

        For compound roles (sub_roles set), returns list[tuple[Shape, Shape]].
        For simple roles, returns list[Shape].
        """
        # Compound role detection: pair_indices / pair_names in find_params
        if role.sub_roles:
            return self._discover_pairs(slide, role)

        strategy = role.find_strategy
        params = role.find_params or {}

        if strategy == ShapeFindStrategy.BY_INDEX:
            return self._by_index(slide, params, role.cardinality)
        if strategy == ShapeFindStrategy.BY_NAME:
            return self._by_name(slide, params)
        if strategy == ShapeFindStrategy.BY_GROUP_CHILDREN:
            return self._by_group_children(slide, params)
        if strategy == ShapeFindStrategy.BY_POSITION:
            return self._by_position(slide, params, role.cardinality)
        return []

    # ------------------------------------------------------------------
    # Strategy implementations
    # ------------------------------------------------------------------

    def _by_index(self, slide, params: dict, cardinality: int) -> list:
        """Lookup by 'index' (single) or 'indices' (list)."""
        if 'indices' in params:
            return [self._safe_get_by_index(slide, i) for i in params['indices']]
        if 'index' in params:
            shape = self._safe_get_by_index(slide, params['index'])
            return [shape] if shape is not None else []
        return []

    def _by_name(self, slide, params: dict) -> list:
        """Recursive search by exact name. Accepts 'names' (list) or 'name' (single)."""
        if 'names' in params:
            return [self._find_by_name_recursive(slide.shapes, n) for n in params['names']]
        if 'name' in params:
            shape = self._find_by_name_recursive(slide.shapes, params['name'])
            return [shape] if shape is not None else []
        return []

    def _by_group_children(self, slide, params: dict) -> list:
        """Find a parent group by name, return children sorted by position."""
        group_name = params.get('group_name')
        if not group_name:
            return []
        group = self._find_by_name_recursive(slide.shapes, group_name)
        if group is None or not hasattr(group, 'shapes'):
            return []
        axis = params.get('sort_axis', 'top')
        children = list(group.shapes)
        return self._sort_by_position(children, axis)

    def _by_position(self, slide, params: dict, cardinality: int) -> list:
        """Filter by shape_type, sort by 'top' or 'left', take first N (or all)."""
        if MSO_SHAPE_TYPE is None:
            return []
        type_filter = params.get('shape_type')
        axis = params.get('sort_axis', 'top')
        candidates = list(slide.shapes)
        if type_filter is not None:
            candidates = [s for s in candidates if getattr(s, 'shape_type', None) == type_filter]
        sorted_shapes = self._sort_by_position(candidates, axis)
        if isinstance(cardinality, int) and cardinality > 0:
            return sorted_shapes[:cardinality]
        return sorted_shapes

    # ------------------------------------------------------------------
    # Compound (pair) discovery for sub_roles
    # ------------------------------------------------------------------

    def _discover_pairs(self, slide, role: ShapeRole) -> List[Tuple[Any, Any]]:
        """Compound roles: return list of (title_shape, body_shape) tuples."""
        params = role.find_params or {}
        pairs: list[tuple] = []

        if 'pair_indices' in params:
            for ti, bi in params['pair_indices']:
                pairs.append((self._safe_get_by_index(slide, ti),
                              self._safe_get_by_index(slide, bi)))
        elif 'pair_names' in params:
            for tn, bn in params['pair_names']:
                pairs.append((self._find_by_name_recursive(slide.shapes, tn),
                              self._find_by_name_recursive(slide.shapes, bn)))
        return pairs

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _safe_get_by_index(slide, idx: int):
        try:
            if 0 <= idx < len(slide.shapes):
                return slide.shapes[idx]
        except Exception:
            return None
        return None

    def _find_by_name_recursive(self, shapes, name: str):
        """Recursive shape search by exact name, traversing GROUP children."""
        if MSO_SHAPE_TYPE is None:
            for s in shapes:
                if getattr(s, 'name', None) == name:
                    return s
            return None
        try:
            for s in shapes:
                if s.name == name:
                    return s
                try:
                    if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                        found = self._find_by_name_recursive(s.shapes, name)
                        if found is not None:
                            return found
                except AttributeError:
                    pass
        except Exception:
            pass
        return None

    @staticmethod
    def _sort_by_position(shapes: list, axis: str) -> list:
        """Sort shapes by 'top' or 'left' coordinate. Missing coords sink to end."""
        key = (lambda s: getattr(s, 'top', 0) or 0) if axis == 'top' \
            else (lambda s: getattr(s, 'left', 0) or 0)
        return sorted(shapes, key=key)

    # ------------------------------------------------------------------
    # Excess-shape clearing
    # ------------------------------------------------------------------

    @staticmethod
    def clear_text(shape) -> None:
        """Clear text content of a shape if it has a text frame."""
        if shape is None:
            return
        try:
            if shape.has_text_frame:
                shape.text_frame.text = ''
        except Exception:
            pass
